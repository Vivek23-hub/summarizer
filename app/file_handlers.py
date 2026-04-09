import pdfplumber
import docx
import pytesseract
from PIL import Image
import tempfile
from pptx import Presentation
import numpy as np
import io


def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            if page.extract_text():
                text += page.extract_text() + "\n"
    return text


def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_from_txt(file):
    return file.read().decode("utf-8")


def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


import pytesseract
from PIL import Image


def extract_text_from_image(file):
    image = Image.open(file)

    image = image.convert("L")
    image = image.resize((image.width * 2, image.height * 2))

    text = pytesseract.image_to_string(image)

    print("\n===== CLEAN OCR TEXT =====\n", text)

    return text


def extract_text(file, filename):
    ext = filename.split(".")[-1].lower()

    if ext == "pdf":
        return extract_text_from_pdf(file)

    elif ext in ["doc", "docx"]:
        return extract_text_from_docx(file)

    elif ext == "txt":
        return extract_text_from_txt(file)

    elif ext in ["ppt", "pptx"]:
        return extract_text_from_pptx(file)

    elif ext in ["png", "jpg", "jpeg"]:
        return extract_text_from_image(file)

    else:
        raise ValueError("Unsupported file type")
